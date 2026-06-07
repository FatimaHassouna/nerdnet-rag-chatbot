import os
import json
import logging
import datetime
import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from llama_index.core import VectorStoreIndex, StorageContext, Document, load_index_from_storage
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI as LLamaOpenAI
import streamlit.components.v1 as components
from PIL import Image
import requests
from io import BytesIO
from pptx import Presentation
import tempfile
import fitz  # PyMuPDF
import base64
import sqlite3
import hashlib
import streamlit as st
from dotenv import load_dotenv
from openai import OpenAI
from llama_index.core import VectorStoreIndex, StorageContext, Document, load_index_from_storage
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.llms.openai import OpenAI as LLamaOpenAI
import streamlit.components.v1 as components
from PIL import Image
import requests
import base64
from typing import Optional, Tuple
import pandas as pd  # <-- Add this line




AUTH_DB = "nerdnet_auth.db"
# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

def test_sidebar():
    st.sidebar.button("📊 Test Button", on_click=lambda: st.toast("Button works!"))
    st.write("Button should appear in sidebar")

if __name__ == "__main__":
    test_sidebar()  # Temporarily replace main() call
# Constants
PERSIST_DIR = "./storage"
CHAT_HISTORY_FILE = "chat_history.json"
PDF_PATH = "networkbook.pdf"
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


#### UI & Display Functions ######################################################################3

def load_css():
    st.markdown("""
    <style>
    .chat-message {
        padding: 12px;
        border-radius: 8px;
        margin-bottom: 8px;
    }
    .analytics-card {
        padding: 15px;
        border-radius: 10px;
        background-color: #f8f9fa;
        margin-bottom: 15px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    .metric-value {
        font-size: 24px;
        font-weight: bold;
        color: #1E88E5;
    }
    .user-message {
        background-color: #f0f2f6;
        border-left: 4px solid #4e79a7;
    }
    .assistant-message {
        background-color: #e6f3ff;
        border-left: 4px solid #2e7d32;
    }
    .sidebar .sidebar-content {
        background-color: #f8f9fa;
    }
    .stButton button {
        border-radius: 20px;
        width: 100%;
        margin-bottom: 5px;
    }
    .element-container button:hover {
        background-color: #dfe6f3 !important;
    }
    .listen-btn {
        background-color: #f0f2f6;
        border: 1px solid #4e79a7;
        padding: 4px 8px;
        border-radius: 12px;
        font-size: 12px;
        margin-top: 8px;
        display: inline-flex;
        align-items: center;
    }
    .stop-btn {
        background-color: #ffebee;
        border: 1px solid #f44336;
        padding: 4px 8px;
        border-radius: 12px;
        font-size: 12px;
        margin-top: 8px;
        display: inline-flex;
        align-items: center;
    }
    .button-container {
        display: flex;
        gap: 8px;
        margin-top: 8px;
    }
    .stSidebar .stButton button {
        border-radius: 0px !important;
        width: 100%;
        margin-bottom: 8px;
    }
    .image-container {
        margin-top: 12px;
        margin-bottom: 12px;
        border: 1px solid #e0e0e0;
        border-radius: 8px;
        padding: 8px;
    }
    .pptx-download-btn {
        margin-top: 12px;
        margin-bottom: 12px;
    }
    .chat-actions {
        display: flex;
        gap: 8px;
        margin-top: 8px;
        margin-bottom: 8px;
    }
    .chat-actions button {
        padding: 6px 12px;
        font-size: 14px;
        border-radius: 12px;
        border: 1px solid #e0e0e0;
        background-color: white;
    }
    .chat-actions button:hover {
        background-color: #f5f5f5 !important;
    }
    .copy-btn {
        border-color: #4CAF50 !important;
        color: #4CAF50 !important;
    }
    .edit-btn {
        border-color: #2196F3 !important;
        color: #2196F3 !important;
    }
    .save-btn {
        border-color: #4CAF50 !important;
        background-color: #4CAF50 !important;
        color: white !important;
    }
    .cancel-btn {
        border-color: #f44336 !important;
        background-color: #f44336 !important;
        color: white !important;
    }
    </style>
    """, unsafe_allow_html=True)


# Date grouping
def categorize_date(iso_timestamp):
    today = datetime.date.today()
    try:
        date = datetime.date.fromisoformat(iso_timestamp.split("T")[0])
    except:
        return "Unknown"
    delta = (today - date).days
    if delta == 0:
        return "Today"
    elif delta == 1:
        return "Yesterday"
    elif delta <= 7:
        return "Last 7 Days"
    else:
        return date.strftime("%b %d, %Y")

def show_analytics_dashboard():
    try:
        import pandas as pd  # Double safety import
        
        st.title("📊 Analytics Dashboard")
        if not st.session_state.chats:
            st.warning("No chat data available yet!")
            return
        try:
            all_messages = []
            for chat_name, messages in st.session_state.chats.items():
                for msg in messages:
                    if msg["role"] != "system":
                        all_messages.append({
                            "date": datetime.datetime.fromisoformat(msg["timestamp"]).date(),
                            "chat": chat_name,
                            "role": msg["role"],
                            "length": len(msg["content"])
                        })
            
            df = pd.DataFrame(all_messages)   
        except Exception as e:
            st.error(f"Data processing error: {str(e)}")
            
    except ImportError:
        st.error("Please install pandas: pip install pandas")



def show_analytics_dashboard():
    st.title("📊 Advanced Analytics Dashboard")
    
    if not st.session_state.chats:
        st.warning("No chat data available yet!")
        return
        
    try:
        # Process all messages into a DataFrame
        all_messages = []
        for chat_name, messages in st.session_state.chats.items():
            for msg in messages:
                if msg["role"] != "system":
                    try:
                        dt = datetime.datetime.fromisoformat(msg["timestamp"])
                        all_messages.append({
                            "date": dt.date(),
                            "datetime": dt,
                            "chat": chat_name,
                            "role": msg["role"],
                            "length": len(msg["content"]),
                            "is_user": msg["role"] == "user"
                        })
                    except:
                        continue
        
        if not all_messages:
            st.warning("No message data available for analytics")
            return
            
        df = pd.DataFrame(all_messages)
        
        # Calculate metrics
        total_chats = len(st.session_state.chats)
        total_messages = len(df)
        user_messages = df[df['is_user']].shape[0]
        assistant_messages = total_messages - user_messages
        avg_message_length = df['length'].mean()
        
        # Display metrics
        col1, col2, col3 = st.columns(3)
        with col1:
            st.markdown("""
            <div class="analytics-card">
                <h3>Total Chats</h3>
                <p class="metric-value">{}</p>
            </div>
            """.format(total_chats), unsafe_allow_html=True)
            
        with col2:
            st.markdown("""
            <div class="analytics-card">
                <h3>Total Messages</h3>
                <p class="metric-value">{}</p>
            </div>
            """.format(total_messages), unsafe_allow_html=True)
            
        with col3:
            st.markdown("""
            <div class="analytics-card">
                <h3>Avg Message Length</h3>
                <p class="metric-value">{:.1f}</p>
            </div>
            """.format(avg_message_length), unsafe_allow_html=True)
        
        # Message type distribution
        st.subheader("Message Distribution")
        chart_data = pd.DataFrame({
            'Type': ['User Messages', 'Assistant Messages'],
            'Count': [user_messages, assistant_messages]
        })
        st.bar_chart(chart_data.set_index('Type'))
        
        # Message activity over time
        st.subheader("Message Activity Over Time")
        daily_counts = df.groupby(df['datetime'].dt.date).size()
        st.line_chart(daily_counts)
        
    except Exception as e:
        st.error(f"Error generating analytics: {str(e)}")



######### Session State & Chat Control Functions #######################################################################################
# Session state initialization
def init_session_state():
    if "current_chat" not in st.session_state:
        st.session_state.current_chat = "default"
    if "chats" not in st.session_state:
        st.session_state.chats = {"default": [{
            "role": "system",
            "content": "Chat started.",
            "timestamp": datetime.datetime.now().isoformat()
        }]}
    if "pending_prompt" not in st.session_state:
        st.session_state.pending_prompt = None
    if "awaiting_response" not in st.session_state:
        st.session_state.awaiting_response = False
    if "index" not in st.session_state:
        st.session_state.index = None
    if "is_speaking" not in st.session_state:
        st.session_state.is_speaking = False
    if "auto_read" not in st.session_state:
        st.session_state.auto_read = True
    if "clipboard" not in st.session_state:
        st.session_state.clipboard = ""
    if "editing" not in st.session_state:
        st.session_state.editing = None
    if "edit_content" not in st.session_state:
        st.session_state.edit_content = ""
def delete_chat(chat_name: str):
    """Delete chat from database"""
    if "user_id" in st.session_state:
        delete_user_chat(st.session_state["user_id"], chat_name)
        del st.session_state.chats[chat_name]
        st.session_state.current_chat = next(iter(st.session_state.chats.keys()))
        st.rerun()

def rename_chat(old_name: str, new_name: str):
    """Rename chat in database"""
    if "user_id" in st.session_state and old_name in st.session_state.chats:
        chat_data = st.session_state.chats[old_name]
        delete_user_chat(st.session_state["user_id"], old_name)
        st.session_state.chats[new_name] = chat_data
        del st.session_state.chats[old_name]
        save_user_chat(st.session_state["user_id"], new_name, chat_data)
        if st.session_state.current_chat == old_name:
            st.session_state.current_chat = new_name
        st.rerun()


# Load chat history
def load_chat_history():
    try:
        if os.path.exists(CHAT_HISTORY_FILE):
            with open(CHAT_HISTORY_FILE, "r") as f:
                chats = json.load(f)
                if "default" not in chats:
                    chats["default"] = [{
                        "role": "system",
                        "content": "Chat started.",
                        "timestamp": datetime.datetime.now().isoformat()
                    }]
                return chats
    except Exception as e:
        logger.error(f"Error loading chat history: {str(e)}")
    return {"default": [{
        "role": "system",
        "content": "Chat started.",
        "timestamp": datetime.datetime.now().isoformat()
    }]}

def save_chat_history():
    """Save chat history to database"""
    if "user_id" in st.session_state and "current_chat" in st.session_state:
        save_user_chat(
            st.session_state["user_id"],
            st.session_state["current_chat"],
            st.session_state.chats[st.session_state["current_chat"]]
        )



#############################################################################################################
 # AI $ Multimedia

# Model init
def initialize_components():
    try:
        embed_model = OpenAIEmbedding(model="text-embedding-3-large", dimensions=1536)
        llm = LLamaOpenAI(model="gpt-4", streaming=True)
        return embed_model, llm
    except Exception as e:
        st.error(f"Failed to initialize models: {str(e)}")
        logger.error(f"Model initialization error: {str(e)}")
        return None, None

def get_or_create_index(embed_model, llm):
    try:
        if st.session_state.index is not None:
            return st.session_state.index

        if not os.path.exists(PERSIST_DIR):
            with st.spinner("Creating new index from textbook..."):
                documents = []
                with fitz.open(PDF_PATH) as doc:
                    for i, page in enumerate(doc):
                        text = page.get_text().strip()
                        if text:
                            documents.append(Document(
                                text=text,
                                metadata={
                                    "page_number": i + 1,
                                    "source": "Computer Networking: A Top-Down Approach"
                                }
                            ))
                            
                index = VectorStoreIndex.from_documents(documents, embed_model=embed_model, llm=llm)
                index.storage_context.persist(persist_dir=PERSIST_DIR)
        else:
            with st.spinner("Loading existing index from storage..."):       
                storage_context = StorageContext.from_defaults(persist_dir=PERSIST_DIR)
                index = load_index_from_storage(storage_context, embed_model=embed_model)

        st.session_state.index = index
        return index

    except Exception as e:
        st.error(f"Index error: {str(e)}")
        logger.error(f"Index error: {str(e)}")
        return None

def query_with_sources(index, prompt, embed_model, llm):
    try:
        query_engine = index.as_query_engine(llm=llm, embed_model=embed_model, return_source=True)
        response = query_engine.query(prompt)

        # Extract and format sources
        source_texts = []
        for node in response.source_nodes:
            page = node.metadata.get("page_number", "Unknown")
            src = node.metadata.get("source", "Unknown source")
            source_texts.append(f"- Page {page} from {src}")

        full_response = str(response.response) + "\n\nSources:\n" + "\n".join(source_texts)
        return full_response

    except Exception as e:
        st.error(f"Query failed: {str(e)}")
        logger.error(f"Query error: {str(e)}")
        return "An error occurred during query."

def stream_query_with_sources(index, prompt, embed_model, llm, chat_id):
    try:
        # Initialize streaming response
        placeholder = st.empty()
        full_response = ""
        
        # Get streaming response from OpenAI directly
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            stream=True
        )
        
        # Stream the response
        for chunk in response:
            if chunk.choices[0].delta.content is not None:
                full_response += chunk.choices[0].delta.content
                placeholder.markdown(f"🤖 **Assistant:** {full_response}")

        # Now query index (non-streaming) to get sources
        query_engine = index.as_query_engine(llm=llm, embed_model=embed_model, return_source=True)
        response = query_engine.query(prompt)

        sources = []
        for node in response.source_nodes:
            page = node.metadata.get("page_number", "Unknown")
            src = node.metadata.get("source", "Unknown source")
            sources.append(f"- Page {page} from {src}")

        full_response = full_response + "\n\nSources:\n" + "\n".join(sources)

        # Save full response to chat history
        st.session_state.chats[chat_id].append({
            "role": "assistant",
            "content": full_response,
            "timestamp": datetime.datetime.now().isoformat()
        })

        return full_response

    except Exception as e:
        st.error(f"Streaming query failed: {str(e)}")
        logger.error(f"Streaming query error: {str(e)}")
        return "Streaming error occurred."

def init_session_state():
    # First load from file if exists
    if "chats" not in st.session_state:
        st.session_state.chats = load_chat_history()
    
    # Then set defaults only if empty
    if not st.session_state.chats:
        st.session_state.chats = {
            "default": [{
                "role": "system",
                "content": "Chat started.",
                "timestamp": datetime.datetime.now().isoformat()
            }]
        }
    
    if "current_chat" not in st.session_state:
        st.session_state.current_chat = "default"
    
    # Initialize other states
    if "pending_prompt" not in st.session_state:
        st.session_state.pending_prompt = None
    if "awaiting_response" not in st.session_state:
        st.session_state.awaiting_response = False
    if "index" not in st.session_state:
        st.session_state.index = None
    if "is_speaking" not in st.session_state:
        st.session_state.is_speaking = False
    if "auto_read" not in st.session_state:
        st.session_state.auto_read = True
    if "clipboard" not in st.session_state:
        st.session_state.clipboard = ""
    if "editing" not in st.session_state:
        st.session_state.editing = None
    if "edit_content" not in st.session_state:
        st.session_state.edit_content = ""
    if "show_analytics" not in st.session_state:
        st.session_state.show_analytics = False



# Function to trigger browser TTS with better handling
def speak_text(text):
    safe_text = json.dumps(text)[1:-1]  # Escape quotes and special characters
    components.html(f"""
    <script>
    function speak() {{
        if (!window.speechSynthesis) {{
            console.error("Speech synthesis not supported");
            return;
        }}
        
        // Cancel any ongoing speech
        window.speechSynthesis.cancel();
        
        const msg = new SpeechSynthesisUtterance();
        msg.text = {json.dumps(text)};
        msg.volume = 1;
        msg.rate = 1;
        msg.pitch = 1;
        
        msg.onend = function() {{
            console.log("Speech synthesis finished");
        }};
        
        msg.onerror = function(event) {{
            console.error("Speech synthesis error", event);
        }};
        
        window.speechSynthesis.speak(msg);
    }}
    
    // Check if voices are loaded
    if (window.speechSynthesis.getVoices().length > 0) {{
        speak();
    }} else {{
        window.speechSynthesis.onvoiceschanged = speak;
    }}
    </script>
    """, height=0)
def stop_speech():
    components.html("""
    <script>
    if (window.speechSynthesis) {
        window.speechSynthesis.cancel();
        console.log("Speech stopped");
    }gpt
    </script>
    """, height=0)

# Image generation function
def generate_image(prompt):
    try:
        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
        )
        image_url = response.data[0].url
        return image_url
    except Exception as e:
        st.error(f"Image generation failed: {str(e)}")
        return None

def get_base64_image(image_path):
    with open(image_path, "rb") as img_file:
        encoded = base64.b64encode(img_file.read()).decode()
        return f"data:image/png;base64,{encoded}"

def interactive_speak(text):
    idle_img = get_base64_image("image.png")
    talking_img = get_base64_image("image2.gif")
    safe_text = json.dumps(text)[1:-1]

    components.html(f"""
    <div id="avatar-container" style="text-align:center; margin-top:10px;">
        <img id="avatar-img" src="{idle_img}" width="150" />
        <div class="button-container">
            <button onclick="startSpeaking()" class="listen-btn">🔊 Speak</button>
            <button onclick="stopSpeaking()" class="stop-btn">⏹️ Stop</button>
        </div>
        <p id="speech-status" style="font-size:12px; color:#555;">Idle</p>
    </div>

    <script>
    const avatarImg = document.getElementById("avatar-img");
    const statusText = document.getElementById("speech-status");

    const msg = new SpeechSynthesisUtterance();
    msg.text = "{safe_text}";
    msg.volume = 1;
    msg.rate = 1;
    msg.pitch = 1;

    msg.onstart = function() {{
        avatarImg.src = "{talking_img}";
        statusText.innerText = "Speaking...";
    }};

    msg.onend = function() {{
        avatarImg.src = "{idle_img}";
        statusText.innerText = "Finished.";
    }};

    msg.onerror = function() {{
        avatarImg.src = "{idle_img}";
        statusText.innerText = "Error.";
    }};

    function startSpeaking() {{
        window.speechSynthesis.cancel();
        window.speechSynthesis.speak(msg);
    }}

    function stopSpeaking() {{
        window.speechSynthesis.cancel();
        avatarImg.src = "{idle_img}";
        statusText.innerText = "Stopped.";
    }}
    </script>
    """, height=300)

# PowerPoint generation functions
def create_powerpoint_from_content(content):
    """Convert text content into a PowerPoint presentation"""
    try:
        prs = Presentation()
        
        # Split content into slides (assuming each slide starts with "Slide X:")
        slides_content = [slide.strip() for slide in content.split("Slide ") if slide.strip()]
        
        for slide_content in slides_content:
            if not slide_content:
                continue
                
            # Get slide number and actual content
            slide_parts = slide_content.split(":", 1)
            if len(slide_parts) < 2:
                continue
                
            content = slide_parts[1].strip()
            
            # Add a slide with title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title (slide number)
            title = slide.shapes.title
            title.text = f"Slide {slide_parts[0]}"
            
            # Add content
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content
            
        # Save to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            tmp.seek(0)
            pptx_bytes = tmp.read()
            
        return pptx_bytes
    except Exception as e:
        st.error(f"Failed to create PowerPoint: {str(e)}")
        return None

def save_powerpoint(pptx_bytes, filename="presentation.pptx"):
    """Create a download button for the PowerPoint"""
    if pptx_bytes:
        st.download_button(
            label="📥 Download PowerPoint",
            data=pptx_bytes,
            file_name=filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"ppt_{filename}_{datetime.datetime.now().timestamp()}",
            use_container_width=True,
            help="Click to download the PowerPoint presentation"
        )


#################################################################################################################

# Authentication & User Management Functions




# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def init_db():
    """Initialize the SQLite database for user authentication"""
    conn = sqlite3.connect(AUTH_DB)
    c = conn.cursor()
    
    # Drop tables if they exist (for clean start)
    c.execute('DROP TABLE IF EXISTS users')
    c.execute('DROP TABLE IF EXISTS user_chats')
    
    # Create users table
    c.execute('''
        CREATE TABLE users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            email TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Create user_chats table
    c.execute('''
        CREATE TABLE user_chats (
            user_id INTEGER NOT NULL,
            chat_name TEXT NOT NULL,
            chat_data TEXT NOT NULL,
            last_modified TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            PRIMARY KEY (user_id, chat_name),
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
    ''')
    
    # Create indexes
    c.execute('CREATE INDEX idx_username ON users(username)')
    c.execute('CREATE INDEX idx_email ON users(email)')
    
    conn.commit()
    conn.close()


def create_user(username: str, password: str, email: Optional[str] = None) -> bool:
    """Create a new user in the database"""
    try:
        # First check if username exists
        conn = sqlite3.connect(AUTH_DB)
        c = conn.cursor()
        c.execute('SELECT 1 FROM users WHERE username = ?', (username,))
        if c.fetchone():
            st.error("Username already exists")
            return False
            
        # If email provided, check if email exists
        if email:
            c.execute('SELECT 1 FROM users WHERE email = ?', (email,))
            if c.fetchone():
                st.error("Email already exists")
                return False
                
        # If checks pass, create the user
        password_hash = hashlib.sha256(password.encode()).hexdigest()
        c.execute('INSERT INTO users (username, password_hash, email) VALUES (?, ?, ?)',
                 (username, password_hash, email))
        conn.commit()
        return True
    except Exception as e:
        logger.error(f"Error creating user: {str(e)}")
        return False
    finally:
        conn.close()

def verify_user(username: str, password: str) -> bool:
    """Verify user credentials"""
    password_hash = hashlib.sha256(password.encode()).hexdigest()
    conn = sqlite3.connect(AUTH_DB)
    c = conn.cursor()
    c.execute('SELECT password_hash FROM users WHERE username = ?', (username,))
    result = c.fetchone()
    conn.close()
    return result and result[0] == password_hash

def get_user_id(username: str) -> Optional[int]:
    """Get user ID from username"""
    conn = sqlite3.connect(AUTH_DB)
    c = conn.cursor()
    c.execute('SELECT id FROM users WHERE username = ?', (username,))
    result = c.fetchone()
    conn.close()
    return result[0] if result else None

def save_user_chat(user_id: int, chat_name: str, chat_data: dict):
    """Save chat data for a specific user"""
    conn = sqlite3.connect(AUTH_DB)
    c = conn.cursor()
    chat_json = json.dumps(chat_data)
    try:
        c.execute('''
            INSERT INTO user_chats (user_id, chat_name, chat_data)
            VALUES (?, ?, ?)
            ON CONFLICT(user_id, chat_name) DO UPDATE SET
            chat_data = excluded.chat_data,
            last_modified = CURRENT_TIMESTAMP
        ''', (user_id, chat_name, chat_json))
        conn.commit()
    finally:
        conn.close()

def load_user_chats(user_id: int) -> dict:
    """Load all chats for a specific user"""
    conn = sqlite3.connect(AUTH_DB)
    c = conn.cursor()
    c.execute('SELECT chat_name, chat_data FROM user_chats WHERE user_id = ?', (user_id,))
    chats = {row[0]: json.loads(row[1]) for row in c.fetchall()}
    conn.close()
    return chats or {"Default Chat": [{
        "role": "system",
        "content": "New chat started",
        "timestamp": datetime.datetime.now().isoformat()
    }]}

def delete_user_chat(user_id: int, chat_name: str):
    """Delete a specific chat for a user"""
    conn = sqlite3.connect(AUTH_DB)
    c = conn.cursor()
    c.execute('DELETE FROM user_chats WHERE user_id = ? AND chat_name = ?', (user_id, chat_name))
    conn.commit()
    conn.close()

# Initialize database
init_db()

# # Authentication UI
def show_auth_page():
    """Display authentication page with login and signup tabs"""
    st.title("🔒 NerdNET Authentication")
    
    tab1, tab2 = st.tabs(["Login", "Sign Up"])
    
    with tab1:
        with st.form("login_form"):
            username = st.text_input("Username", key="login_username")
            password = st.text_input("Password", type="password", key="login_password")
            submitted = st.form_submit_button("Login")
            
            if submitted:
                if verify_user(username, password):
                    st.session_state["authenticated"] = True
                    st.session_state["username"] = username
                    st.session_state["user_id"] = get_user_id(username)
                    # Initialize user chats
                    st.session_state.chats = load_user_chats(st.session_state["user_id"])
                    st.session_state.current_chat = next(iter(st.session_state.chats.keys()), "Default Chat")
                    st.rerun()
                else:
                    st.error("Invalid username or password")
    with tab2:
        with st.form("signup_form"):
            username = st.text_input("Username", key="signup_username")
            email = st.text_input("Email (optional)", key="signup_email")
            password = st.text_input("Password", type="password", key="signup_password")
            confirm_password = st.text_input("Confirm Password", type="password", key="confirm_password")
            submitted = st.form_submit_button("Sign Up")
            
            if submitted:
                if password != confirm_password:
                    st.error("Passwords don't match")
                elif len(username) < 3:
                    st.error("Username must be at least 3 characters")
                elif len(password) < 6:
                    st.error("Password must be at least 6 characters")
                else:
                    if create_user(username, password, email):
                        st.success("Account created successfully! Please login.")
                    # Error messages are now handled within create_user()


